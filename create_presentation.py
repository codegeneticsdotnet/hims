from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

def add_slide(title_text, content_items):
    # Use the 'Title and Content' layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = title_text
    
    # Set content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hospital Information Management System"
subtitle.text = "Process Flow & System Architecture Analysis"

# Slide 2: High Level Overview
add_slide("High Level Process Flow", [
    "1. Patient Registration (Entry Point)",
    "2. Clinical Assessment (OPD/IPD)",
    "3. Care & Treatment (Nursing/Meds)",
    "4. Ancillary Services (Lab/Pharmacy)",
    "5. Billing & Finance (Exit Point)"
])

# Slide 3: Patient Registration
add_slide("1. Patient Registration", [
    "Input: Patient Demographics (Name, Age, Address)",
    "Action: Search existing or Create new record",
    "Output: Unique Patient ID (PID)",
    "Key Table: patient_personal_info",
    "Roles Involved: Receptionist, Clinic Staff"
])

# Slide 4: Out-Patient (OPD) Flow
add_slide("2. Out-Patient (OPD) Flow", [
    "Visit Type: Walk-in or Appointment",
    "Action: Doctor Consultation",
    "Clinical Data: Vitals, Diagnosis, Complaints",
    "Orders: Prescriptions, Lab Tests",
    "Disposition: Home (Discharge) or Admit (to IPD)",
    "Key Tables: patient_details_iop, iop_diagnosis"
])

# Slide 5: In-Patient (IPD) Flow
add_slide("3. In-Patient (IPD) Flow", [
    "Admission: Room/Ward Assignment (Room Transfer)",
    "Care: Daily Nurse Notes, Vital Signs Monitoring",
    "Procedures: Bedside procedures (Oxygen, Nebulizer)",
    "Transfers: Tracking movement for billing accuracy",
    "Key Tables: iop_room_transfer, iop_nurse_notes, iop_progress_note"
])

# Slide 6: Ancillary Services
add_slide("4. Ancillary Services", [
    "Laboratory:",
    "  - Service Request -> Result Entry -> Verification",
    "  - Table: lab_service_request",
    "Pharmacy:",
    "  - Prescription -> Stock Check -> Dispense",
    "  - Inventory: Real-time deduction (pharmacy_sales)",
    "  - Table: pharmacy_inventory_details"
])

# Slide 7: Billing & Discharge Logic
add_slide("5. Billing & Discharge", [
    "Charge Aggregation Logic:",
    "  - Room: Days x Room Rate (from transfers)",
    "  - Doctor: Professional Fee (from Discharge Advice)",
    "  - Items: Meds + Labs + Services",
    "Payment: Cash/Card -> Receipt Generation",
    "Discharge: Patient Status updated to 'Discharged'",
    "Key Tables: iop_billing, iop_receipt"
])

# Slide 8: Database Analysis - Unused Tables
add_slide("Database Optimization: Unused Tables", [
    "Legacy Inventory Tables (Prefix 'o_'):",
    "  - o_batch, o_items, o_inventory, o_transactions",
    "Empty Functional Tables (Candidates for Cleanup):",
    "  - cart_billing, category_list, declaredor",
    "  - doctors_fee (Replaced by iop_discharge_advice)",
    "  - patient_appointment (If unused)",
    "  - surgical_package"
])

# Save the presentation
output_file = 'HIMS_Process_Flow.pptx'
prs.save(output_file)
print(f"Presentation saved to {output_file}")
